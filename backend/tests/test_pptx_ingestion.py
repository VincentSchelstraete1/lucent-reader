import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.ingestion import DocumentExtractionError, PptxDocumentIngestor
from app.ingestion.pptx_ingestor import PPTX_MIME_TYPE
from app.main import app
from app.normalization import normalize_document
from app.routers.ingestion import get_pptx_ingestor


def _tiny_png() -> bytes:
    import struct
    import zlib

    def chunk(tag: bytes, data: bytes) -> bytes:
        body = tag + data
        return struct.pack(">I", len(data)) + body + struct.pack(">I", zlib.crc32(body))

    signature = b"\x89PNG\r\n\x1a\n"
    header = struct.pack(">IIBBBBB", 3, 3, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + b"\x00\xff\x00" * 3 for _ in range(3))
    data = zlib.compress(raw)
    return signature + chunk(b"IHDR", header) + chunk(b"IDAT", data) + chunk(b"IEND", b"")


def _sample_pptx(*, include_image: bool = False) -> bytes:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Cache Associativity"
    body = first.placeholders[1]
    body.text_frame.text = "Direct-mapped caches place a block in exactly one location."
    first.notes_slide.notes_text_frame.text = "Remind students about the tradeoff with set-associative caches."

    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = second.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = "Level"
    table.cell(0, 1).text = "Latency"
    table.cell(1, 0).text = "L1"
    table.cell(1, 1).text = "4 cycles"
    if include_image:
        second.shapes.add_picture(io.BytesIO(_tiny_png()), Inches(1), Inches(3))

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class StubPptxIngestor:
    def __init__(self, *, error: bool = False) -> None:
        self.error = error
        self.calls: list[tuple[bytes, str]] = []

    def ingest_pptx(self, pptx_bytes: bytes, *, filename: str):
        self.calls.append((pptx_bytes, filename))
        if self.error:
            raise DocumentExtractionError("pptx extractor detail that must not leak")
        from app.ingestion import RawContentBlock, RawDocument, RawPage, SourceLocation

        block = RawContentBlock(
            id="slide-1-shape-1",
            page_number=None,
            type="text",
            bbox=None,
            reading_order=0,
            text="stub slide text",
            location=SourceLocation(kind="slide", index=1, sequence_id="shape-1"),
        )
        page = RawPage(
            None, "stub slide text", [block], location=SourceLocation(kind="slide", index=1)
        )
        return RawDocument(
            source_type="pptx",
            filename=filename,
            page_count=1,
            markdown="stub slide text",
            pages=[page],
            images=[],
        )


@pytest.fixture()
def stub_pptx_ingestor():
    stub = StubPptxIngestor()
    app.dependency_overrides[get_pptx_ingestor] = lambda: stub
    yield stub
    app.dependency_overrides.pop(get_pptx_ingestor, None)


def test_pptx_reaches_normalized_document_with_multiple_slides():
    raw = PptxDocumentIngestor().ingest_pptx(_sample_pptx(), filename="deck.pptx")
    assert raw.source_type == "pptx"
    assert raw.page_count == 2
    assert [page.page_number for page in raw.pages] == [None, None]
    assert [page.location.index for page in raw.pages] == [1, 2]
    assert all(page.location.kind == "slide" for page in raw.pages)

    normalized = normalize_document(raw)
    assert len(normalized.pages) == 2
    first_types = [block.type for block in normalized.pages[0].blocks]
    assert "heading" in first_types
    assert "paragraph" in first_types
    second_types = [block.type for block in normalized.pages[1].blocks]
    assert "table" in second_types


def test_pptx_slide_number_never_written_to_legacy_page_number_field():
    raw = PptxDocumentIngestor().ingest_pptx(_sample_pptx(), filename="deck.pptx")
    for page in raw.pages:
        assert page.page_number is None
        for block in page.blocks:
            assert block.page_number is None
            assert block.location.kind == "slide"
            assert block.location.index == page.location.index


def test_pptx_speaker_notes_are_captured_separately_from_visible_shapes():
    raw = PptxDocumentIngestor().ingest_pptx(_sample_pptx(), filename="deck.pptx")
    notes_blocks = [block for block in raw.pages[0].blocks if block.location.sequence_id == "notes"]
    assert len(notes_blocks) == 1
    assert "set-associative" in (notes_blocks[0].text or "")


def test_pptx_image_is_extracted_with_slide_provenance_and_bbox():
    raw = PptxDocumentIngestor().ingest_pptx(_sample_pptx(include_image=True), filename="deck.pptx")
    assert len(raw.images) == 1
    image = raw.images[0]
    assert image.page_number is None
    assert image.location.kind == "slide"
    assert image.location.index == 2
    assert image.width == 3
    assert image.height == 3
    assert image.mime_type == "image/png"
    assert image.bbox is not None

    normalized = normalize_document(raw)
    assert len(normalized.images) == 1
    assert normalized.images[0].location.kind == "slide"
    assert normalized.images[0].location.index == 2


def test_pptx_extraction_failure_raises_document_extraction_error():
    with pytest.raises(DocumentExtractionError):
        PptxDocumentIngestor().ingest_pptx(b"not a real pptx file", filename="broken.pptx")


def test_pptx_upload_endpoint_returns_normalized_document(client, stub_pptx_ingestor):
    response = client.post(
        "/ingestion/pptx",
        files={"file": ("deck.pptx", b"PK\x03\x04stub pptx bytes", PPTX_MIME_TYPE)},
    )
    assert response.status_code == 200
    result = response.json()
    assert result["source_type"] == "pptx"
    assert result["pages"][0]["page_number"] is None
    assert result["pages"][0]["location"]["kind"] == "slide"
    assert result["pages"][0]["location"]["index"] == 1
    assert result["normalized"]["pages"][0]["page_number"] is None
    assert stub_pptx_ingestor.calls == [(b"PK\x03\x04stub pptx bytes", "deck.pptx")]


def test_pptx_upload_rejects_wrong_media_type(client, stub_pptx_ingestor):
    response = client.post("/ingestion/pptx", files={"file": ("deck.txt", b"not pptx", "text/plain")})
    assert response.status_code == 415
    assert response.json()["detail"]["code"] == "unsupported_file_type"
    assert stub_pptx_ingestor.calls == []


def test_pptx_upload_rejects_invalid_zip_signature(client, stub_pptx_ingestor):
    response = client.post(
        "/ingestion/pptx",
        files={"file": ("deck.pptx", b"not a zip file", PPTX_MIME_TYPE)},
    )
    assert response.status_code == 422
    assert response.json()["detail"]["code"] == "invalid_pptx"


def test_pptx_upload_requires_authentication(unauthenticated_client, stub_pptx_ingestor):
    response = unauthenticated_client.post(
        "/ingestion/pptx",
        files={"file": ("deck.pptx", b"PK\x03\x04stub pptx bytes", PPTX_MIME_TYPE)},
    )
    assert response.status_code == 401
    assert stub_pptx_ingestor.calls == []


def test_pptx_extraction_failure_returns_safe_structured_error(client):
    app.dependency_overrides[get_pptx_ingestor] = lambda: StubPptxIngestor(error=True)
    try:
        response = client.post(
            "/ingestion/pptx",
            files={"file": ("broken.pptx", b"PK\x03\x04stub pptx bytes", PPTX_MIME_TYPE)},
        )
    finally:
        app.dependency_overrides.pop(get_pptx_ingestor, None)
    assert response.status_code == 422
    assert response.json()["detail"]["code"] == "pptx_extraction_failed"
    assert "extractor detail" not in response.text


@pytest.mark.integration
def test_real_pptx_extractor_returns_slides_markdown_and_provenance():
    result = PptxDocumentIngestor().ingest_pptx(_sample_pptx(), filename="fixture.pptx")
    assert "Direct-mapped caches" in result.pages[0].text
    assert "Direct-mapped caches" in result.markdown
    assert result.page_count == 2
