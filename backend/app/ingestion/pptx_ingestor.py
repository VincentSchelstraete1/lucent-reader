import base64
import hashlib
from io import BytesIO
from typing import TYPE_CHECKING

from app.ingestion.base import (
    DocumentExtractionError,
    RawContentBlock,
    RawDocument,
    RawImage,
    RawPage,
    SourceLocation,
)
from app.ingestion.markitdown_adapter import MarkItDownAdapter

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape


PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
EMU_PER_POINT = 12700


def _mime_type_from_content_type(content_type: str | None) -> str | None:
    return content_type.split(";")[0].strip() if content_type else None


def _shape_bbox(shape: "BaseShape") -> tuple[float, float, float, float] | None:
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    left, top = shape.left / EMU_PER_POINT, shape.top / EMU_PER_POINT
    return (left, top, left + shape.width / EMU_PER_POINT, top + shape.height / EMU_PER_POINT)


def _table_text(table) -> str:
    rows = [[" ".join(cell.text.split()) for cell in row.cells] for row in table.rows]
    if not rows or not rows[0]:
        return ""
    separator = ["---"] * len(rows[0])
    lines = [rows[0], separator, *rows[1:]]
    return "\n".join("| " + " | ".join(row) + " |" for row in lines)


class PptxDocumentIngestor:
    """Extract PPTX structure into the shared RawDocument model.

    Slide number is meaningful physical provenance (like a PDF page), but it is
    never stored in the legacy page_number field - see SourceLocation's
    docstring. Each RawPage represents one slide, kind="slide".
    """

    def __init__(self, markdown_extractor: MarkItDownAdapter | None = None) -> None:
        self._markdown_extractor = markdown_extractor or MarkItDownAdapter()

    def ingest_pptx(self, pptx_bytes: bytes, *, filename: str) -> RawDocument:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except Exception as exc:
            raise DocumentExtractionError("The PPTX could not be opened") from exc

        try:
            presentation = Presentation(BytesIO(pptx_bytes))
        except Exception as exc:
            raise DocumentExtractionError("The PPTX could not be opened") from exc

        pages: list[RawPage] = []
        images: list[RawImage] = []

        try:
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_location = SourceLocation(kind="slide", index=slide_index)
                blocks: list[RawContentBlock] = []
                slide_images: list[RawImage] = []
                errors: list[str] = []
                text_lines: list[str] = []
                order = 0

                for shape_index, shape in enumerate(slide.shapes, start=1):
                    bbox = _shape_bbox(shape)
                    sequence_id = f"shape-{shape_index}"
                    location = SourceLocation(kind="slide", index=slide_index, sequence_id=sequence_id)
                    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                        text = shape.text_frame.text
                        blocks.append(
                            RawContentBlock(
                                id=f"slide-{slide_index}-{sequence_id}",
                                page_number=None,
                                type="text",
                                bbox=bbox,
                                reading_order=order,
                                text=text,
                                location=location,
                            )
                        )
                        text_lines.append(text)
                        order += 1
                    elif getattr(shape, "has_table", False):
                        table_text = _table_text(shape.table)
                        if table_text:
                            blocks.append(
                                RawContentBlock(
                                    id=f"slide-{slide_index}-{sequence_id}",
                                    page_number=None,
                                    type="text",
                                    bbox=bbox,
                                    reading_order=order,
                                    text=table_text,
                                    location=location,
                                )
                            )
                            text_lines.append(table_text)
                            order += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            if not image_bytes:
                                raise ValueError("empty image blob")
                            width, height = image.size
                            mime_type = _mime_type_from_content_type(image.content_type)
                        except Exception:
                            errors.append(f"Picture shape {shape_index} did not contain usable raster data")
                            continue
                        digest = hashlib.sha256(image_bytes).hexdigest()[:16]
                        image_id = f"slide-{slide_index}-image-{shape_index}-{digest}"
                        encoded = base64.b64encode(image_bytes).decode("ascii")
                        media_type = mime_type or "application/octet-stream"
                        image_id_block = RawContentBlock(
                            id=f"slide-{slide_index}-{sequence_id}",
                            page_number=None,
                            type="image",
                            bbox=bbox,
                            reading_order=order,
                            image_id=image_id,
                            location=location,
                        )
                        blocks.append(image_id_block)
                        slide_images.append(
                            RawImage(
                                id=image_id,
                                page_number=None,
                                bbox=bbox,
                                width=width,
                                height=height,
                                mime_type=mime_type,
                                asset_reference=f"data:{media_type};base64,{encoded}",
                                location=location,
                            )
                        )
                        order += 1

                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        blocks.append(
                            RawContentBlock(
                                id=f"slide-{slide_index}-notes",
                                page_number=None,
                                type="text",
                                bbox=None,
                                reading_order=order,
                                text=notes_text,
                                location=SourceLocation(kind="slide", index=slide_index, sequence_id="notes"),
                            )
                        )
                        text_lines.append(notes_text)

                pages.append(
                    RawPage(
                        page_number=None,
                        text="\n\n".join(text_lines),
                        blocks=blocks,
                        extraction_errors=errors,
                        location=slide_location,
                    )
                )
                images.extend(slide_images)
        except DocumentExtractionError:
            raise
        except Exception as exc:
            raise DocumentExtractionError("The PPTX could not be read") from exc

        markdown = self._markdown_extractor.extract_markdown(
            BytesIO(pptx_bytes),
            filename=filename,
            mimetype=PPTX_MIME_TYPE,
            extension=".pptx",
        )

        return RawDocument(
            source_type="pptx",
            filename=filename,
            page_count=len(pages),
            markdown=markdown,
            pages=pages,
            images=images,
            extraction_metadata={
                "page_extractor": "python-pptx",
                "markdown_extractor": "MarkItDown",
                "slide_count": len(pages),
            },
        )
